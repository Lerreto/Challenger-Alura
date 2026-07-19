from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from nebula_rag.errors import DocumentProcessingError
from nebula_rag.loaders import extract_sections, parse_markdown


def test_markdown_preserves_frontmatter_and_heading_metadata() -> None:
    content = b"""---
document_id: NT-FAQ-001
title: Preguntas frecuentes
category: Atencion
version: '1.0'
owner: Equipo de soporte
language: es-CO
---
# FAQ
## Envios
El envio tarda dos dias.
"""

    metadata, sections = parse_markdown(content.decode())

    assert metadata["document_id"] == "NT-FAQ-001"
    assert metadata["owner"] == "Equipo de soporte"
    assert sections[0].location == "FAQ > Envios"
    assert "El envio tarda" in sections[0].text


def test_heading_immediately_followed_by_a_subheading_is_not_a_content_stub() -> None:
    # "# FAQ" carries no text of its own before "## Envios": indexing it as a
    # standalone chunk would only echo the title and outscore real content on
    # queries that mention the document by name.
    _, sections = parse_markdown("# FAQ\n## Envios\nEl envio tarda dos dias.\n")
    assert len(sections) == 1
    assert sections[0].location == "FAQ > Envios"


def test_heading_with_its_own_text_before_a_subheading_is_kept() -> None:
    _, sections = parse_markdown(
        "# FAQ\nPreguntas frecuentes de la tienda.\n## Envios\nEl envio tarda dos dias.\n"
    )
    assert len(sections) == 2
    assert sections[0].location == "FAQ"
    assert sections[0].text == "FAQ\nPreguntas frecuentes de la tienda."
    assert sections[1].location == "FAQ > Envios"


def test_heading_only_document_falls_back_to_the_raw_body_instead_of_empty() -> None:
    # Every heading here is content-free, but parse_markdown still must not
    # return an empty list — extract_sections treats that as empty_document.
    _, sections = parse_markdown("# Solo título\n## Subtítulo vacío\n")
    assert len(sections) == 1
    assert sections[0].location == "Documento"


def test_representative_uploaded_formats_keep_locations(tmp_path: Path) -> None:
    text_file = tmp_path / "notes.txt"
    text_file.write_text("Linea uno\nLinea dos", encoding="utf-8")

    docx_file = tmp_path / "manual.docx"
    doc = Document()
    doc.add_heading("Garantias", level=1)
    doc.add_paragraph("Cobertura por doce meses")
    doc.save(docx_file)

    xlsx_file = tmp_path / "catalog.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Productos"
    sheet.append(["Producto", "Garantia"])
    sheet.append(["Laptop", "12 meses"])
    workbook.save(xlsx_file)

    pptx_file = tmp_path / "training.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Devoluciones"
    slide.placeholders[1].text = "Solicitar en 15 dias"
    presentation.save(pptx_file)

    assert extract_sections(text_file)[0].location == "Texto"
    assert extract_sections(docx_file)[0].location == "Garantias"
    assert extract_sections(xlsx_file)[0].location == "Hoja: Productos"
    assert extract_sections(pptx_file)[0].location == "Diapositiva 1"


@pytest.mark.parametrize(
    ("name", "content", "needle"),
    [
        ("table.csv", b"producto,precio\nmouse,50000\n", "mouse"),
        ("data.json", b'{"politica": {"dias": 15}}', "politica"),
        ("page.html", b"<h1>Privacidad</h1><p>No vendemos datos.</p>", "No vendemos"),
        ("readme.md", b"# Terminos\nSolo Colombia", "Solo Colombia"),
    ],
)
def test_structured_text_formats_are_supported(
    tmp_path: Path, name: str, content: bytes, needle: str
) -> None:
    path = tmp_path / name
    path.write_bytes(content)
    sections = extract_sections(path)
    assert needle in "\n".join(section.text for section in sections)


def test_empty_pdf_requires_ocr(tmp_path: Path) -> None:
    path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with path.open("wb") as handle:
        writer.write(handle)

    with pytest.raises(DocumentProcessingError, match="ocr_required"):
        extract_sections(path)


def test_zip_bomb_guard_rejects_absurd_office_archive(tmp_path: Path) -> None:
    path = tmp_path / "bomb.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", "x" * 250_000)

    with pytest.raises(DocumentProcessingError, match="extracted_content_too_large"):
        extract_sections(path, max_extracted_bytes=50_000)


def test_archive_member_limit_is_enforced_before_office_parser(tmp_path: Path) -> None:
    path = tmp_path / "many.docx"
    with zipfile.ZipFile(path, "w") as archive:
        for index in range(4):
            archive.writestr(f"word/item-{index}.xml", "safe")

    with pytest.raises(DocumentProcessingError, match="too_many_archive_members"):
        extract_sections(path, max_archive_members=3)


def test_json_depth_and_node_limits_are_enforced(tmp_path: Path) -> None:
    path = tmp_path / "deep.json"
    path.write_text('{"a":{"b":{"c":{"d":1}}}}', encoding="utf-8")

    with pytest.raises(DocumentProcessingError, match="json_too_complex"):
        extract_sections(path, max_json_depth=2)


def test_frontmatter_rejects_deep_yaml_before_parser_recursion() -> None:
    nested = "\n".join(f"{'  ' * level}level-{level}:" for level in range(800))
    content = f"---\n{nested}\n---\n# Documento\nTexto"

    with pytest.raises(DocumentProcessingError, match="invalid_frontmatter"):
        parse_markdown(content)


@pytest.mark.parametrize(
    "frontmatter",
    [
        "owner:\n  name: Legal",
        "tags:\n  - legal",
        "alias: &owner Legal\ncopy: *owner",
        "value: !custom tagged",
    ],
)
def test_frontmatter_accepts_only_flat_scalar_mappings(frontmatter: str) -> None:
    with pytest.raises(DocumentProcessingError, match="invalid_frontmatter"):
        parse_markdown(f"---\n{frontmatter}\n---\n# Documento\nTexto")
