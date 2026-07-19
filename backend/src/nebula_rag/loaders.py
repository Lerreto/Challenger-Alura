from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from pathlib import Path
from typing import Any

import yaml
from bs4 import BeautifulSoup
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from .domain import ExtractedSection
from .errors import DocumentProcessingError


SUPPORTED_EXTENSIONS = {
    ".md",
    ".txt",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".csv",
    ".json",
    ".html",
    ".htm",
}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
MAX_FRONTMATTER_BYTES = 64 * 1024
MAX_FRONTMATTER_LINES = 256
MAX_FRONTMATTER_FIELDS = 128


def _clean(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _read_utf8(path: Path, *, signature: bool = False) -> str:
    try:
        return path.read_text(encoding="utf-8-sig" if signature else "utf-8")
    except UnicodeDecodeError as exc:
        raise DocumentProcessingError("invalid_text_encoding") from exc


def _check_archive(path: Path, max_bytes: int, max_members: int = 2_000) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
            if len(members) > max_members:
                raise DocumentProcessingError("too_many_archive_members")
            total = sum(member.file_size for member in members)
            if total > max_bytes:
                raise DocumentProcessingError("extracted_content_too_large: archivo expandido excede el limite")
            for member in members:
                if member.compress_size and member.file_size / member.compress_size > 250:
                    raise DocumentProcessingError("extracted_content_too_large: compresion sospechosa")
    except zipfile.BadZipFile as exc:
        raise DocumentProcessingError("invalid_office_archive") from exc


def _validate_flat_frontmatter(source: str) -> None:
    try:
        encoded_size = len(source.encode("utf-8"))
    except UnicodeError as exc:
        raise DocumentProcessingError("invalid_frontmatter") from exc
    lines = source.splitlines()
    if encoded_size > MAX_FRONTMATTER_BYTES or len(lines) > MAX_FRONTMATTER_LINES:
        raise DocumentProcessingError("invalid_frontmatter")

    fields = 0
    for line in lines:
        if not line.strip() or line.lstrip().startswith("#"):
            continue
        if line != line.lstrip() or "\t" in line:
            raise DocumentProcessingError("invalid_frontmatter")
        match = re.fullmatch(r"[A-Za-z_][A-Za-z0-9_-]{0,63}\s*:\s*.*", line)
        if not match or re.search(r"(^|\s)[&*!][A-Za-z_]", line):
            raise DocumentProcessingError("invalid_frontmatter")
        fields += 1
        if fields > MAX_FRONTMATTER_FIELDS:
            raise DocumentProcessingError("invalid_frontmatter")


def parse_markdown(text: str) -> tuple[dict[str, Any], list[ExtractedSection]]:
    metadata: dict[str, Any] = {}
    body = text
    if text.startswith("---\n"):
        marker = text.find("\n---\n", 4)
        if marker >= 0:
            frontmatter = text[4:marker]
            _validate_flat_frontmatter(frontmatter)
            try:
                loaded = yaml.safe_load(frontmatter) or {}
            except (yaml.YAMLError, RecursionError, TypeError, UnicodeError) as exc:
                raise DocumentProcessingError("invalid_frontmatter") from exc
            if not isinstance(loaded, dict):
                raise DocumentProcessingError("invalid_frontmatter")
            if any(
                not isinstance(key, str)
                or not isinstance(value, (str, int, float, bool, type(None)))
                for key, value in loaded.items()
            ):
                raise DocumentProcessingError("invalid_frontmatter")
            metadata = loaded
            body = text[marker + 5 :]

    sections: list[ExtractedSection] = []
    headings: list[str] = []
    current: list[str] = []
    location = "Documento"
    # True right after a heading, until a real body line appears under it.
    heading_only = False

    def flush() -> None:
        cleaned = _clean("\n".join(current))
        # A heading immediately followed by another heading (no free text of
        # its own) carries no information beyond its title, which is already
        # exposed via `location` on every other chunk — skip it so it can't
        # outscore chunks with real content on literal title-name queries.
        if cleaned and not heading_only:
            sections.append(ExtractedSection(cleaned, location))

    for line in body.splitlines():
        match = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if not match:
            if line.strip():
                heading_only = False
            current.append(line)
            continue
        flush()
        current = []
        level = len(match.group(1))
        title = match.group(2).strip().strip("#").strip()
        headings = headings[: level - 1]
        headings.append(title)
        location = " > ".join(headings)
        current.append(title)
        heading_only = True
    flush()
    if not sections and _clean(body):
        sections = [ExtractedSection(_clean(body), "Documento")]
    return metadata, sections


def _extract_pdf(path: Path) -> list[ExtractedSection]:
    try:
        reader = PdfReader(path)
        sections = []
        for index, page in enumerate(reader.pages, 1):
            text = _clean(page.extract_text() or "")
            if text:
                sections.append(ExtractedSection(text, f"Página {index}", {"page": index}))
    except Exception as exc:
        raise DocumentProcessingError("invalid_pdf") from exc
    if not sections:
        raise DocumentProcessingError(
            "ocr_required: el PDF no contiene texto seleccionable; OCR no esta disponible"
        )
    return sections


def _extract_docx(path: Path) -> list[ExtractedSection]:
    document = WordDocument(path)
    sections: list[ExtractedSection] = []
    current: list[str] = []
    location = "Documento"

    def flush() -> None:
        text = _clean("\n".join(current))
        if text:
            sections.append(ExtractedSection(text, location))

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style and paragraph.style.name.lower().startswith("heading"):
            flush()
            current = [text]
            location = text
        else:
            current.append(text)
    flush()
    for table_index, table in enumerate(document.tables, 1):
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        text = _clean("\n".join(rows))
        if text:
            sections.append(ExtractedSection(text, f"Tabla {table_index}"))
    return sections


def _extract_xlsx(path: Path) -> list[ExtractedSection]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(value or f"Columna {index + 1}") for index, value in enumerate(rows[0])]
        rendered = []
        for row in rows[1:]:
            values = [f"{headers[index]}: {value}" for index, value in enumerate(row) if value is not None]
            if values:
                rendered.append(" | ".join(values))
        if not rendered:
            rendered = [" | ".join(headers)]
        sections.append(
            ExtractedSection(
                _clean("\n".join(rendered)),
                f"Hoja: {sheet.title}",
                {"sheet": sheet.title},
            )
        )
    workbook.close()
    return sections


def _extract_pptx(path: Path) -> list[ExtractedSection]:
    presentation = Presentation(path)
    sections = []
    for index, slide in enumerate(presentation.slides, 1):
        fragments: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                fragments.append(shape.text.strip())
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    fragments.append(f"Notas: {shape.text.strip()}")
        text = _clean("\n".join(fragments))
        if text:
            sections.append(
                ExtractedSection(text, f"Diapositiva {index}", {"slide": index})
            )
    return sections


def _extract_csv(path: Path) -> list[ExtractedSection]:
    text = _read_utf8(path, signature=True)
    rows = list(csv.DictReader(io.StringIO(text)))
    rendered = [" | ".join(f"{key}: {value}" for key, value in row.items()) for row in rows]
    return [ExtractedSection(_clean("\n".join(rendered) or text), "Tabla CSV")]


def _render_json(value: Any, prefix: str = "") -> list[str]:
    if isinstance(value, dict):
        lines: list[str] = []
        for key, child in value.items():
            location = f"{prefix}.{key}" if prefix else str(key)
            lines.extend(_render_json(child, location))
        return lines
    if isinstance(value, list):
        lines = []
        for index, child in enumerate(value):
            lines.extend(_render_json(child, f"{prefix}[{index}]"))
        return lines
    return [f"{prefix}: {value}"]


def _validate_json_complexity(value: Any, max_depth: int, max_nodes: int) -> None:
    nodes = 0
    stack: list[tuple[Any, int]] = [(value, 0)]
    while stack:
        current, depth = stack.pop()
        nodes += 1
        if nodes > max_nodes or depth > max_depth:
            raise DocumentProcessingError("json_too_complex")
        if isinstance(current, dict):
            stack.extend((child, depth + 1) for child in current.values())
        elif isinstance(current, list):
            stack.extend((child, depth + 1) for child in current)


def _extract_json(
    path: Path, *, max_depth: int, max_nodes: int
) -> list[ExtractedSection]:
    try:
        value = json.loads(_read_utf8(path))
    except json.JSONDecodeError as exc:
        raise DocumentProcessingError("invalid_json") from exc
    _validate_json_complexity(value, max_depth, max_nodes)
    return [ExtractedSection(_clean("\n".join(_render_json(value))), "Datos JSON")]


def _extract_html(path: Path) -> list[ExtractedSection]:
    soup = BeautifulSoup(_read_utf8(path), "html.parser")
    for unwanted in soup(["script", "style", "noscript"]):
        unwanted.decompose()
    text = _clean(soup.get_text("\n"))
    return [ExtractedSection(text, "Página HTML")]


def extract_sections(
    path: Path,
    *,
    max_extracted_bytes: int = 4 * 1024 * 1024,
    max_archive_members: int = 2_000,
    max_json_depth: int = 32,
    max_json_nodes: int = 100_000,
) -> list[ExtractedSection]:
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise DocumentProcessingError("unsupported_extension")
    if extension in OFFICE_EXTENSIONS:
        _check_archive(path, max_extracted_bytes, max_archive_members)
    elif path.stat().st_size > max_extracted_bytes:
        raise DocumentProcessingError("extracted_content_too_large")

    try:
        if extension == ".md":
            _, sections = parse_markdown(_read_utf8(path))
        elif extension == ".txt":
            sections = [ExtractedSection(_clean(_read_utf8(path)), "Texto")]
        elif extension == ".pdf":
            sections = _extract_pdf(path)
        elif extension == ".docx":
            sections = _extract_docx(path)
        elif extension == ".xlsx":
            sections = _extract_xlsx(path)
        elif extension == ".pptx":
            sections = _extract_pptx(path)
        elif extension == ".csv":
            sections = _extract_csv(path)
        elif extension == ".json":
            sections = _extract_json(
                path, max_depth=max_json_depth, max_nodes=max_json_nodes
            )
        else:
            sections = _extract_html(path)
    except DocumentProcessingError:
        raise
    except Exception as exc:
        raise DocumentProcessingError("invalid_document") from exc

    total = sum(len(section.text.encode("utf-8")) for section in sections)
    if total > max_extracted_bytes:
        raise DocumentProcessingError("extracted_content_too_large")
    if not sections or not any(section.text.strip() for section in sections):
        raise DocumentProcessingError("empty_document")
    return sections


def read_frontmatter(path: Path) -> dict[str, Any]:
    if path.suffix.lower() != ".md":
        return {}
    metadata, _ = parse_markdown(_read_utf8(path))
    return metadata
